# Simple text extraction helpers for common file types

from __future__ import annotations

from typing import Optional
from pathlib import Path

from pypdf import PdfReader
try:
    from docx import Document as DocxDocument  # type: ignore
except Exception:  # pragma: no cover - optional dependency in test env
    DocxDocument = None  # type: ignore
try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover - optional dependency in test env
    Presentation = None  # type: ignore


def extract_text(path: str | Path) -> str:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(p)
    if suffix in (".docx", ".doc"):
        return _extract_docx(p)
    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(p)
    # Fallback for txt and others
    try:
        return p.read_text(errors="ignore")
    except Exception:
        return ""


def _extract_pdf(p: Path) -> str:
    try:
        reader = PdfReader(str(p))
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        return "\n".join(texts)
    except Exception:
        return ""


def _extract_docx(p: Path) -> str:
    try:
        if DocxDocument is None:
            return ""
        doc = DocxDocument(str(p))
        return "\n".join(par.text for par in doc.paragraphs)
    except Exception:
        return ""


def _extract_pptx(p: Path) -> str:
    try:
        if Presentation is None:
            return ""
        prs = Presentation(str(p))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception:
        return ""
